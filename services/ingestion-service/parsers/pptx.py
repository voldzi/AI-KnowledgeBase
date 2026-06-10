from __future__ import annotations

import io
import re

from app.object_storage import SourceObject
from parsers.base import DocumentParser, ParsedBlock, ParserError, ParserResult, ParserUnavailable

PPTX_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


class PptxParser(DocumentParser):
    name = "pptx"

    def supports(self, source: SourceObject) -> bool:
        filename = source.filename.lower()
        return source.mime_type in PPTX_MIME_TYPES or filename.endswith(".pptx")

    def parse(self, source: SourceObject, *, parser_profile: str) -> ParserResult:
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover
            raise ParserUnavailable("PPTX_PARSER_UNAVAILABLE", "python-pptx is not installed") from exc

        try:
            presentation = Presentation(io.BytesIO(source.content))
        except Exception as exc:
            raise ParserError("PPTX_PARSE_FAILED", f"Presentation could not be opened: {exc.__class__.__name__}") from exc

        blocks: list[ParsedBlock] = []
        warnings: list[tuple[str, str]] = []
        tables_detected = 0
        offset = 0

        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_title = _slide_title(slide) or f"Slide {slide_number}"
            section_path = [slide_title]
            slide_texts: list[tuple[str, str]] = []  # (block_type, text)

            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    tables_detected += 1
                    rows = []
                    for row in shape.table.rows:
                        cells = [_clean(cell.text) for cell in row.cells]
                        line = " | ".join(cell for cell in cells if cell)
                        if line:
                            rows.append(line)
                    if rows:
                        slide_texts.append(("table", "\n".join(rows)))
                    continue
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = _clean(shape.text_frame.text)
                if not text:
                    continue
                block_type = "heading" if text == slide_title else "paragraph"
                slide_texts.append((block_type, text))

            notes = _speaker_notes(slide)
            if notes:
                slide_texts.append(("paragraph", f"Poznámky: {notes}"))

            for block_type, text in slide_texts:
                blocks.append(
                    ParsedBlock(
                        text=text,
                        page_number=slide_number,
                        section_path=list(section_path),
                        section_title=slide_title,
                        article_number=None,
                        paragraph_number=None,
                        char_start=offset,
                        char_end=offset + len(text),
                        block_type=block_type,
                    )
                )
                offset += len(text) + 2

        if not blocks:
            warnings.append(("NO_TEXT_EXTRACTED", "Presentation contains no readable text."))

        return ParserResult(
            parser_name=self.name,
            blocks=blocks,
            pages_processed=len(list(presentation.slides)),
            tables_detected=tables_detected,
            warnings=warnings,
        )


def _slide_title(slide: object) -> str | None:
    shapes = getattr(slide, "shapes", None)
    title = getattr(shapes, "title", None) if shapes is not None else None
    if title is not None and getattr(title, "has_text_frame", False):
        text = _clean(title.text_frame.text)
        return text or None
    return None


def _speaker_notes(slide: object) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""
    notes_slide = slide.notes_slide
    frame = getattr(notes_slide, "notes_text_frame", None)
    if frame is None:
        return ""
    return _clean(frame.text)


def _clean(value: str) -> str:
    return re.sub(r"\s+", " ", value or "").strip()
