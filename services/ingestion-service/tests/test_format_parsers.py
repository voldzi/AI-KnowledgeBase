from __future__ import annotations

import hashlib
import io

from openpyxl import Workbook

from app.object_storage import SourceObject
from parsers.html import HtmlParser
from parsers.text import TextParser
from parsers.xlsx import XlsxParser


def _source(filename: str, mime_type: str, content: bytes) -> SourceObject:
    return SourceObject(
        uri=f"s3://test/{filename}",
        filename=filename,
        mime_type=mime_type,
        content=content,
        sha256=f"sha256:{hashlib.sha256(content).hexdigest()}",
    )


def test_html_parser_extracts_headings_and_paragraphs() -> None:
    markup = """
    <html><head><title>Ignored</title><style>p { color: red; }</style></head>
    <body>
      <h1>Směrnice pro řízení dokumentů</h1>
      <p>Výjimku ze směrnice schvaluje gestor dokumentu.</p>
      <h2>Platnost</h2>
      <p>Platná verze je verze se stavem valid.</p>
      <table><tr><td>Role</td><td>Odpovědnost</td></tr></table>
      <script>alert("skip me");</script>
    </body></html>
    """.encode("utf-8")
    parser = HtmlParser()
    source = _source("smernice.html", "text/html", markup)

    assert parser.supports(source)
    result = parser.parse(source, parser_profile="default")

    texts = [block.text for block in result.blocks]
    assert "Směrnice pro řízení dokumentů" in texts
    assert any("gestor dokumentu" in text for text in texts)
    assert not any("alert" in text for text in texts)
    assert not any("color: red" in text for text in texts)
    assert result.tables_detected == 1

    heading_blocks = [block for block in result.blocks if block.block_type == "heading"]
    assert heading_blocks[0].section_path == ["Směrnice pro řízení dokumentů"]
    platnost = next(block for block in result.blocks if block.text == "Platná verze je verze se stavem valid.")
    assert platnost.section_path == ["Směrnice pro řízení dokumentů", "Platnost"]
    assert platnost.section_title == "Platnost"


def test_html_parser_reports_empty_document() -> None:
    parser = HtmlParser()
    result = parser.parse(_source("empty.html", "text/html", b"<html><body></body></html>"), parser_profile="default")
    assert result.blocks == []
    assert ("NO_TEXT_EXTRACTED", "Parser did not extract readable text.") in result.warnings


def test_text_parser_supports_csv_json_and_xml_sources() -> None:
    parser = TextParser()
    cases = [
        _source("export.csv", "text/csv", b"role,owner\nGestor,Legal"),
        _source("payload.json", "application/json", b'{"role":"owner","decision":"approve"}'),
        _source("feed.xml", "application/xml", b"<root><role>auditor</role></root>"),
    ]

    for source in cases:
        assert parser.supports(source)
        result = parser.parse(source, parser_profile="default")
        assert result.blocks


def test_xlsx_parser_extracts_sheet_rows_as_table_blocks() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Rizika"
    sheet.append(["Riziko", "Dopad", "Opatření"])
    sheet.append(["Výpadek LLM Gateway", "Vysoký", "Failover na záložní runtime"])
    sheet.append(["Nedostupnost Qdrant", "Vysoký", "Replikace kolekce"])
    second = workbook.create_sheet("Prázdný")
    buffer = io.BytesIO()
    workbook.save(buffer)

    parser = XlsxParser()
    source = _source(
        "rizika.xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        buffer.getvalue(),
    )

    assert parser.supports(source)
    result = parser.parse(source, parser_profile="default")

    assert result.tables_detected == 1
    assert result.pages_processed == 2
    assert len(result.blocks) == 1
    block = result.blocks[0]
    assert block.block_type == "table"
    assert block.section_path == ["Rizika"]
    assert "Riziko | Dopad | Opatření" in block.text
    assert "Výpadek LLM Gateway | Vysoký | Failover na záložní runtime" in block.text


def test_xlsx_parser_repeats_header_in_continuation_blocks() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["ID", "Hodnota"])
    for index in range(100):
        sheet.append([index, f"radek_{index}"])
    buffer = io.BytesIO()
    workbook.save(buffer)

    parser = XlsxParser()
    result = parser.parse(
        _source("data.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", buffer.getvalue()),
        parser_profile="default",
    )

    assert len(result.blocks) > 1
    for block in result.blocks[1:]:
        assert block.text.splitlines()[0] == "ID | Hodnota"


def test_xlsx_blank_cells_do_not_shift_a_value_into_another_column() -> None:
    workbook = Workbook()
    workbook.active.append(["Service", "Minimum", "Proposed", "Guaranteed"])
    workbook.active.append(["AKB", None, "32 GB", None])
    buffer = io.BytesIO()
    workbook.save(buffer)
    result = XlsxParser().parse(_source("sizing.xlsx", "application/octet-stream", buffer.getvalue()), parser_profile="default")
    assert "AKB |  | 32 GB | " in result.blocks[0].text


def test_docx_retains_body_order_headings_and_table_associations() -> None:
    from docx import Document
    from parsers.docx import DocxParser

    document = Document()
    document.add_heading("Pilot infrastructure", level=1)
    document.add_heading("Sizing proposal", level=2)
    document.add_paragraph("These values are proposals, not measured minima.")
    table = document.add_table(rows=2, cols=3)
    for cell, text in zip(table.rows[0].cells, ["Service", "Minimum", "Proposed"]):
        cell.text = text
    table.rows[1].cells[0].text = "AKB"
    table.rows[1].cells[2].text = "32 GB"
    document.add_heading("Recovery", level=2)
    document.add_paragraph("RTO must be agreed with the recipient.")
    buffer = io.BytesIO()
    document.save(buffer)
    result = DocxParser().parse(_source("pilot.docx", "application/octet-stream", buffer.getvalue()), parser_profile="default")
    texts = [block.text for block in result.blocks]
    sizing = next(block for block in result.blocks if block.block_type == "table")
    assert texts.index(sizing.text) < texts.index("Recovery")
    assert sizing.section_path == ["Pilot infrastructure", "Sizing proposal"]
    assert "AKB |  | 32 GB" in sizing.text
    assert result.blocks[-1].section_path == ["Pilot infrastructure", "Recovery"]
    assert all(block.page_number is None for block in result.blocks)
    assert result.tables_detected == 1


def test_docx_review_warning_survives_good_text_extraction_quality() -> None:
    from docx import Document
    from docx.oxml import OxmlElement

    from app.pipeline import _quality_report
    from parsers.docx import DocxParser

    document = Document()
    paragraph = document.add_paragraph("Readable body text with sufficient density. " * 10)
    paragraph._p.append(OxmlElement("w:ins"))
    buffer = io.BytesIO()
    document.save(buffer)
    parsed = DocxParser().parse(_source("revision.docx", "application/octet-stream", buffer.getvalue()), parser_profile="default")
    quality = _quality_report(parsed, extraction_profile="document_text_v1")
    assert quality.quality_tier == "good"
    assert quality.pages_processed == quality.pages_with_text == 0
    assert quality.requires_review is True
    assert any(code == "DOCX_TRACKED_CHANGES_REQUIRE_REVIEW" for code, _ in parsed.warnings)


def test_pptx_parser_extracts_slides_with_titles_and_notes() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    from parsers.pptx import PptxParser

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Architektura platformy"
    slide.placeholders[1].text = "Platforma se skládá z dokumentového registru a vyhledávání."
    slide.notes_slide.notes_text_frame.text = "Zmínit governance proces."

    table_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    rows, cols = 2, 2
    table_shape = table_slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(1))
    table_shape.table.cell(0, 0).text = "Role"
    table_shape.table.cell(0, 1).text = "Odpovědnost"
    table_shape.table.cell(1, 0).text = "Gestor"
    table_shape.table.cell(1, 1).text = "Schvaluje výjimky"

    buffer = io.BytesIO()
    presentation.save(buffer)

    parser = PptxParser()
    source = _source(
        "architektura.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        buffer.getvalue(),
    )

    assert parser.supports(source)
    result = parser.parse(source, parser_profile="default")

    assert result.pages_processed == 2
    assert result.tables_detected == 1

    slide1_blocks = [block for block in result.blocks if block.page_number == 1]
    assert any(block.block_type == "heading" and block.text == "Architektura platformy" for block in slide1_blocks)
    assert any("dokumentového registru" in block.text for block in slide1_blocks)
    assert any(block.text.startswith("Poznámky:") for block in slide1_blocks)
    assert all(block.section_path == ["Architektura platformy"] for block in slide1_blocks)

    table_blocks = [block for block in result.blocks if block.block_type == "table"]
    assert len(table_blocks) == 1
    assert "Gestor | Schvaluje výjimky" in table_blocks[0].text
